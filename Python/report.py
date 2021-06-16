import os
import sys
import configparser
import glob
import csv
from pptx import Presentation
import pptx
import console


def start():
    print("********************PowerPointの画像適用を開始します。********************")
    Display = console.display()

    # 変数定義
    PYTHON_DIR = os.path.dirname(os.path.abspath(__file__))
    CONFIG_INI_PATH = os.path.join(PYTHON_DIR, "config.ini")

    # iniファイル存在チェック
    if not os.path.exists(CONFIG_INI_PATH):
        msg_list = list()
        msg_list.append("【ERROR】code00001")
        msg_list.append("iniファイルが存在しません。")
        msg_list.append("report.pyと同じ階層にconfig.iniを作成してください。")
        msg_list.append("config.iniは以下のように作成してください。")
        msg_list.append("1.セクション[IMAGE]を追加します。")
        msg_list.append(
            "2.セクション[IMAGE]のパラメータBeforeFolderにトリミング前の画像フォルダを指定します。")
        msg_list.append("3.セクション[IMAGE]のパラメータAfterFolderにトリミング後の画像フォルダを指定します。")
        msg_list.append(
            "4.セクション[IMAGE]のパラメータFileExtensionにトリミングする画像ファイルの拡張子を空白区切りで指定します。")
        msg_list.append("5.セクション[REPORT]を追加します。")
        msg_list.append(
            "6.セクション[REPORT]のパラメータReplaceSettingsFolderに置換設定csvフォルダを指定します。")
        msg_list.append(
            "7.セクション[REPORT]のパラメータBeforeFolderに置換前のPowerPointフォルダを指定します。")
        msg_list.append(
            "8.セクション[REPORT]のパラメータAfterFolderに置換後のPowerPointフォルダを指定します。")
        msg_list.append("(注意点)")
        msg_list.append("トリミング処理の画像フォルダのファイルは再帰的に取得されます。")
        msg_list.append(
            "トリミング前の画像フォルダにフォルダが存在する場合、トリミング後の画像フォルダにフォルダを新たに作成します。")
        msg_list.append("OpenCVで処理できない画像ファイルの拡張子を指定した場合はエラーが発生します。")

        msg_list.append("PowerPoint置換処理の置換設定csvフォルダのファイルは再帰的に取得されます。")
        msg_list.append("置換設定csvフォルダと同じ階層に置換前のPowerPointが無いとエラーが発生します。")
        msg_list.append(
            "置換前のPowerPointフォルダにフォルダが存在する場合、置換後のPowerPointフォルダにフォルダを新たに作成します。")
        Display.error(msg_list)
        print("********************PowerPointの画像適用を終了します。********************")
        sys.exit()

    # iniファイル読み取り
    try:
        config_ini = configparser.ConfigParser()
        config_ini.read(CONFIG_INI_PATH, encoding="utf-8")

        REPLACE_SETTINGS_FOLDER = config_ini["REPORT"]["ReplaceSettingsFolder"]
        BEFORE_FOLDER = config_ini["REPORT"]["BeforeFolder"]
        AFTER_FOLDER = config_ini["REPORT"]["AfterFolder"]
    except KeyError:
        msg_list = list()
        msg_list.append("【ERROR】code00002")
        msg_list.append("iniファイルの中身が存在しません。")
        msg_list.append("config.iniは以下のように作成してください。")
        msg_list.append("1.セクション[IMAGE]を追加します。")
        msg_list.append(
            "2.セクション[IMAGE]のパラメータBeforeFolderにトリミング前の画像フォルダを指定します。")
        msg_list.append("3.セクション[IMAGE]のパラメータAfterFolderにトリミング後の画像フォルダを指定します。")
        msg_list.append(
            "4.セクション[IMAGE]のパラメータFileExtensionにトリミングする画像ファイルの拡張子を空白区切りで指定します。")
        msg_list.append("5.セクション[REPORT]を追加します。")
        msg_list.append(
            "6.セクション[REPORT]のパラメータReplaceSettingsFolderに置換設定csvフォルダを指定します。")
        msg_list.append(
            "7.セクション[REPORT]のパラメータBeforeFolderに置換前のPowerPointフォルダを指定します。")
        msg_list.append(
            "8.セクション[REPORT]のパラメータAfterFolderに置換後のPowerPointフォルダを指定します。")
        msg_list.append("(注意点)")
        msg_list.append("トリミング処理の画像フォルダのファイルは再帰的に取得されます。")
        msg_list.append(
            "トリミング前の画像フォルダにフォルダが存在する場合、トリミング後の画像フォルダにフォルダを新たに作成します。")
        msg_list.append("OpenCVで処理できない画像ファイルの拡張子を指定した場合はエラーが発生します。")

        msg_list.append("PowerPoint置換処理の置換設定csvフォルダのファイルは再帰的に取得されます。")
        msg_list.append("置換設定csvフォルダと同じ階層に置換前のPowerPointが無いとエラーが発生します。")
        msg_list.append(
            "置換前のPowerPointフォルダにフォルダが存在する場合、置換後のPowerPointフォルダにフォルダを新たに作成します。")
        Display.error(msg_list)
        print("********************PowerPointの画像適用を終了します。********************")
        sys.exit()
    except configparser.ParsingError:
        msg_list = list()
        msg_list.append("【ERROR】code00003")
        msg_list.append("iniファイルの書き方が間違っています。")
        msg_list.append("config.iniは以下のように作成してください。")
        msg_list.append("1.セクション[IMAGE]を追加します。")
        msg_list.append(
            "2.セクション[IMAGE]のパラメータBeforeFolderにトリミング前の画像フォルダを指定します。")
        msg_list.append("3.セクション[IMAGE]のパラメータAfterFolderにトリミング後の画像フォルダを指定します。")
        msg_list.append(
            "4.セクション[IMAGE]のパラメータFileExtensionにトリミングする画像ファイルの拡張子を空白区切りで指定します。")
        msg_list.append("5.セクション[REPORT]を追加します。")
        msg_list.append(
            "6.セクション[REPORT]のパラメータReplaceSettingsFolderに置換設定csvフォルダを指定します。")
        msg_list.append(
            "7.セクション[REPORT]のパラメータBeforeFolderに置換前のPowerPointフォルダを指定します。")
        msg_list.append(
            "8.セクション[REPORT]のパラメータAfterFolderに置換後のPowerPointフォルダを指定します。")
        msg_list.append("(注意点)")
        msg_list.append("トリミング処理の画像フォルダのファイルは再帰的に取得されます。")
        msg_list.append(
            "トリミング前の画像フォルダにフォルダが存在する場合、トリミング後の画像フォルダにフォルダを新たに作成します。")
        msg_list.append("OpenCVで処理できない画像ファイルの拡張子を指定した場合はエラーが発生します。")

        msg_list.append("PowerPoint置換処理の置換設定csvフォルダのファイルは再帰的に取得されます。")
        msg_list.append("置換設定csvフォルダと同じ階層に置換前のPowerPointが無いとエラーが発生します。")
        msg_list.append(
            "置換前のPowerPointフォルダにフォルダが存在する場合、置換後のPowerPointフォルダにフォルダを新たに作成します。")
        Display.error(msg_list)
        print("********************PowerPointの画像適用を終了します。********************")
        sys.exit()

    # ファイル一覧取得(再帰的)
    files_path = glob.glob(os.path.join(
        REPLACE_SETTINGS_FOLDER, "**"), recursive=True)
    if len(files_path) < 2:
        msg_list = list()
        msg_list.append("【ERROR】code00004")
        msg_list.append("csvファイルが存在しません。")
        msg_list.append(
            "config.iniのReplaceSettingsFolderに設定したパス配下にcsvファイルを配置してください。")
        msg_list.append("(ReplaceSettingsFolderに設定したパス)")
        msg_list.append(REPLACE_SETTINGS_FOLDER)
        Display.error(msg_list)
        print("********************PowerPointの画像適用を終了します。********************")
        sys.exit()

    for file_path in files_path:
        # 拡張子チェック
        if not os.path.splitext(file_path)[1] == ".csv":
            continue

        # PowerPointファイルパス取得
        before_pptx_dir = os.path.dirname(
            file_path).replace(REPLACE_SETTINGS_FOLDER, BEFORE_FOLDER)
        after_pptx_dir = os.path.dirname(
            file_path).replace(REPLACE_SETTINGS_FOLDER, AFTER_FOLDER)

        pptx_name = os.path.basename(file_path).replace(
            os.path.splitext(file_path)[1], ".pptx"
        )
        before_pptx_path = os.path.join(before_pptx_dir, pptx_name)
        after_pptx_path = os.path.join(after_pptx_dir, pptx_name)
        print("(適用元)" + before_pptx_path)
        print("(適用先)" + after_pptx_path)

        # 画像ファイル一覧の取得
        replace_list = []
        try:
            with open(file_path, "r", encoding="UTF-8", newline="") as csvfile:
                reader = csv.reader(csvfile, delimiter=",", quotechar='"')
                next(reader)

                for row in reader:
                    replace_list.append(row)
        except StopIteration:
            msg_list = list()
            msg_list.append("【ERROR】code00005")
            msg_list.append("csvファイルの中身が存在しません。")
            msg_list.append("csvファイルを以下のように作成してください。")
            msg_list.append("1.ヘッダー行を追加します。")
            msg_list.append("2.2行目以降に" "置換文字列" "," "置換画像ファイルパス" "を追加します。")
            Display.error(msg_list)
            continue

        # PowerPointファイル読み込み
        try:
            prs = Presentation(before_pptx_path)
        except pptx.exc.PackageNotFoundError:
            msg_list = list()
            msg_list.append("【ERROR】code00006")
            msg_list.append("PowerPointファイルが存在しません。")
            msg_list.append("以下の原因が考えられます。")
            msg_list.append("1.拡張子抜きcsvファイル名が拡張子抜きPowerPointファイル名と一致していません。")
            msg_list.append(
                "2.config.iniのBeforeFolderが存在しないフォルダパスになっています。")
            Display.error(msg_list)
            continue

        # PowerPointのフォルダ作成処理

        # 画像の追加と配置元シェイプの削除
        for row_index, replace_row in enumerate(replace_list):
            for i, sld in enumerate(prs.slides, start=1):
                for shape in sld.shapes:
                    # <class 'pptx.shapes.picture.Picture'>を処理するときだけ
                    # shape.textでエラーが発生する
                    try:
                        if shape.text == replace_row[0]:
                            sld.shapes.add_picture(
                                replace_row[1],
                                shape.left,
                                shape.top,
                                width=shape.width,
                                height=shape.height,
                            )
                            sp = shape.element
                            sp.getparent().remove(sp)
                    except AttributeError:
                        continue
                    except FileNotFoundError:
                        msg_list = list()
                        msg_list.append("【ERROR】code00007")
                        msg_list.append("csvファイルに存在しない画像ファイルパスが記載されています。")
                        msg_list.append("(csvファイルパス)")
                        msg_list.append(file_path)
                        msg_list.append("(" + str(row_index + 1) + "行目)")
                        msg_list.append(str(replace_row))
                        msg_list.append("適用処理は続行されます。")
                        Display.error(msg_list)
                        continue

        # フォルダ作成(AfterにBeforeのフォルダが存在しない場合)
        if not os.path.exists(after_pptx_dir):
            os.makedirs(after_pptx_dir, exist_ok=True)
        prs.save(after_pptx_path)

    print("********************PowerPointの画像適用を終了します。********************")
